#!/usr/bin/env python3
"""
TokenMeter 项目演示 PPT 生成脚本
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
PRIMARY_COLOR = RGBColor(102, 126, 234)  # 紫色渐变
SECONDARY_COLOR = RGBColor(118, 75, 162)  # 深紫色
ACCENT_COLOR = RGBColor(78, 205, 196)  # 青色
TEXT_COLOR = RGBColor(51, 51, 51)  # 深灰色
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加背景色块
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(3))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(11.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(28)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """添加章节分隔页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 背景
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = SECONDARY_COLOR
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 标题文字
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 内容列表
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.733), Inches(5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = f"● {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(12)
        p.line_spacing = 1.5
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 标题栏
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = left_box.text_frame
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.3), Inches(5.5), Inches(4.5))
    tf = right_box.text_frame
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.space_before = Pt(8)
    
    return slide

# ==================== 开始生成 PPT ====================

# 第1页：封面
add_title_slide(prs, 
    "TokenMeter", 
    "企业级 MaaS 成本追踪与归因分析平台\n产品方案与开发计划")

# 第2页：项目概述
add_content_slide(prs, "项目概述", [
    "TokenMeter 是一款开源的企业级 AI 成本管理工具",
    "帮助企业追踪、归因和优化多源 MaaS（模型即服务）支出",
    "支持 OpenAI、Azure、Claude、通义千问等主流提供商",
    "提供实时监控、预算预警、成本分摊等核心功能",
    "开源版本免费，企业版 ¥5000/套"
])

# 第3页：市场痛点
add_section_slide(prs, "01 市场痛点")

# 第4页：痛点详情
add_content_slide(prs, "企业 AI 成本管理痛点", [
    "多云分散：同时使用多个 MaaS 提供商，账单分散难统计",
    "黑盒消费：不知道哪个项目、哪个团队花了多少",
    "月底惊吓：只有月底出账单，无法实时监控",
    "维度缺失：云账单只有 API Key 维度，没有业务标签",
    "成本失控：缺乏预算管理，经常超支",
    "分摊困难：无法按项目/团队进行成本分摊"
])

# 第5页：解决方案
add_section_slide(prs, "02 解决方案")

# 第6页：产品功能
add_content_slide(prs, "核心功能", [
    "📊 多源聚合：统一监控 OpenAI、Azure、Claude、通义千问等",
    "🏷️ 精细归因：按项目、团队、环境、用户多维度标记",
    "⚡ 实时监控：分钟级成本更新，告别月底账单惊吓",
    "🚨 智能预警：预算超支、用量突增自动告警",
    "📈 财务级报表：支持成本分摊计算和财务导出",
    "🔒 隐私优先：支持私有化部署，数据不出企业"
])

# 第7页：技术架构
add_section_slide(prs, "03 技术架构")

# 第8页：架构图
add_content_slide(prs, "系统架构", [
    "【接入层】API 代理网关，透明转发，零侵入接入",
    "【处理层】多厂商适配器，统一协议转换",
    "【数据层】SQLite/PostgreSQL 双支持，成本计算引擎",
    "【展示层】Web 仪表盘，实时数据可视化",
    "【通知层】Webhook 集成，飞书/钉钉/Slack 告警",
    "技术栈：Python + FastAPI + SQLAlchemy + Vue.js"
])

# 第9页：支持的厂商
add_two_column_slide(prs, "支持的 MaaS 提供商",
    "国际厂商", [
        "OpenAI (GPT-4/3.5)",
        "Azure OpenAI",
        "Anthropic Claude",
        "Cohere",
        "Mistral"
    ],
    "国内厂商", [
        "阿里云 通义千问",
        "百度 文心一言",
        "智谱 GLM",
        "讯飞星火"
    ]
)

# 第10页：开发进度
add_section_slide(prs, "04 开发进度")

# 第11页：已完成
add_content_slide(prs, "已完成（3天）", [
    "✅ Day 0: 项目启动 - MVP 代码、Demo 上线",
    "✅ Day 1: 基础加固 - 错误处理、日志系统、12个测试",
    "✅ Day 2: 多厂商支持 - 4家厂商适配、17个测试",
    "累计：29个测试用例，100%通过率",
    "代码提交：6 commits",
    "Demo 地址：https://elbert-haustellate-jett.ngrok-free.dev"
])

# 第12页：开发计划
add_content_slide(prs, "开发计划（剩余12天）", [
    "📅 Day 3-4: 预算预警系统（预算管理、飞书通知）",
    "📅 Day 5-6: 用户认证系统（JWT、权限管理）",
    "📅 Day 7-8: 数据导出与备份（CSV、自动备份）",
    "📅 Day 9-10: 性能优化与压测（1000+ RPS）",
    "📅 Day 11-12: 文档完善与演示环境",
    "📅 Day 13-15: 发布推广（GitHub、技术文章）"
])

# 第13页：商业模式
add_section_slide(prs, "05 商业模式")

# 第14页：定价策略
add_two_column_slide(prs, "产品定价",
    "开源版（免费）", [
        "多厂商成本追踪",
        "Web 仪表盘",
        "预算预警",
        "飞书/钉钉通知",
        "数据导出",
        "社区支持"
    ],
    "企业版（¥5000/套）", [
        "SSO/LDAP 集成",
        "高级 RBAC 权限",
        "高级报表与 BI",
        "PostgreSQL 支持",
        "企业级技术支持",
        "定制化开发"
    ]
)

# 第15页：竞争优势
add_content_slide(prs, "竞争优势", [
    "💡 差异化定位：专注财务视角（成本分摊、预算管理）",
    "🌍 中文原生：更好的中文本地化支持",
    "🔧 灵活归因：支持 Project/Team/Env/User 四维标签",
    "🏢 私有化优先：支持完全本地化部署",
    "💰 价格优势：比 Helicone 便宜，比自研省时",
    "📈 开源生态：MIT 协议，社区驱动"
])

# 第16页：竞品分析
add_content_slide(prs, "竞品对比", [
    "Helicone：开源 ✅，但成本归因弱 ❌",
    "LangSmith：功能强 ✅，但闭源收费 ❌",
    "Keywords AI：功能全 ✅，但价格贵 ❌",
    "TokenMeter：开源 ✅ + 财务视角 ✅ + 中文原生 ✅",
    "",
    "核心差异化：更强的财务视角、更好的中文支持、更灵活的归因"
])

# 第17页：目标用户
add_content_slide(prs, "目标用户", [
    "🏢 中大型科技公司（100+员工，月 AI 支出 $1K-$100K）",
    "👨‍💼 CTO/技术 VP - 总体成本把控",
    "👩‍💻 工程经理 - 团队成本管理",
    "💰 FinOps 团队 - 成本分摊与优化",
    "🔧 平台工程师 - 工具维护",
    "🚀 AI 应用创业公司 - 成本转嫁给客户"
])

# 第18页：市场规模
add_content_slide(prs, "市场规模", [
    "📊 2025年中国企业级 MaaS 市场：100-150亿人民币",
    "💵 成本管理工具收费：按消费额的 3-5%",
    "🎯 潜在市场规模（TAM）：3-7.5亿/年",
    "📈 增长趋势：AI 支出快速增长，管理工具严重滞后",
    "🪟 窗口期：现在进入市场，抢占先机"
])

# 第19页：下一步行动
add_section_slide(prs, "06 下一步行动")

# 第20页：行动计划
add_content_slide(prs, "需要合伙人支持", [
    "🧪 种子用户：寻找 3-5 家企业试用 MVP",
    "💼 商务对接：联系潜在客户，收集需求反馈",
    "📢 市场推广：技术社区、社交媒体曝光",
    "💻 开发资源：根据反馈持续迭代产品",
    "🎯 目标：3个月内获得首个付费客户",
    "📅 里程碑：GitHub 100+ stars，5家企业试用"
])

# 第21页：结语
add_title_slide(prs, 
    "让每一分钱都花得明明白白", 
    "TokenMeter - 企业级 AI 成本管理平台\n\n谢谢！")

# 保存 PPT
output_path = "/Users/apple2/.openclaw/workspace/projects/maas-cost-tracker/docs/TokenMeter_产品方案.pptx"
prs.save(output_path)
print(f"✅ PPT 已生成: {output_path}")
print(f"📊 共 {len(prs.slides)} 页")
